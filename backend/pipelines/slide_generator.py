"""
Slide Generator — converts lesson data into branded .pptx files.

One .pptx per lesson, stored in backend/slides/{course_id}/.
Uses a programmatic template: navy header bar, white body, Calibri fonts.
"""

import os
import json
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from pipelines.config import COURSES_FILE

# -------------------------------------------------------
# Brand Constants
# -------------------------------------------------------
BRAND_NAVY    = RGBColor(0x00, 0x31, 0x7A)   # #00317A
BRAND_CYAN    = RGBColor(0x17, 0xBC, 0xE2)   # #17BCE2
BRAND_ORANGE  = RGBColor(0xF7, 0x8F, 0x20)   # #F78F20
BRAND_GREEN   = RGBColor(0x14, 0xC4, 0x96)   # #14C496
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY    = RGBColor(0xEE, 0xEE, 0xEE)
TEXT_BLACK     = RGBColor(0x22, 0x22, 0x22)
BULLET_GRAY   = RGBColor(0x44, 0x44, 0x44)

# Slide dimensions — widescreen 16:9
SLIDE_WIDTH  = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Header bar
HEADER_HEIGHT = Inches(1.3)

# Content area
CONTENT_TOP    = Inches(1.6)
CONTENT_LEFT   = Inches(0.8)
CONTENT_WIDTH  = Inches(11.7)
CONTENT_HEIGHT = Inches(5.4)

# Paths
BASE_DIR   = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
SLIDES_DIR = os.path.join(BASE_DIR, "slides")
LOGO_PATH  = os.path.join(BASE_DIR, "assets", "logo.png")

os.makedirs(SLIDES_DIR, exist_ok=True)


# -------------------------------------------------------
# Helpers — adding shapes with consistent styling
# -------------------------------------------------------

def _add_filled_rect(slide, left, top, width, height, color):
    """Add a solid-colored rectangle (no border)."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def _add_text_box(slide, left, top, width, height, text, font_size,
                  font_color=TEXT_BLACK, bold=False, alignment=PP_ALIGN.LEFT,
                  font_name="Calibri", anchor=MSO_ANCHOR.TOP):
    """Add a text box with a single run of styled text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None  # manual sizing

    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = font_color
    run.font.bold = bold
    run.font.name = font_name
    return txBox


def _add_logo(slide):
    """Add the PhillipCapital logo to top-right if available."""
    if not os.path.exists(LOGO_PATH):
        return
    try:
        logo_width = Inches(2.0)
        logo_left = SLIDE_WIDTH - logo_width - Inches(0.4)
        logo_top = Inches(0.3)
        slide.shapes.add_picture(LOGO_PATH, logo_left, logo_top, width=logo_width)
    except Exception:
        pass  # Non-critical — skip silently


def _add_slide_number(slide, number, total):
    """Small slide-number indicator at bottom-right."""
    text = f"{number} / {total}"
    _add_text_box(
        slide,
        SLIDE_WIDTH - Inches(1.5), SLIDE_HEIGHT - Inches(0.5),
        Inches(1.2), Inches(0.3),
        text, font_size=10, font_color=BRAND_NAVY,
        alignment=PP_ALIGN.RIGHT, bold=False,
    )


# -------------------------------------------------------
# Slide Builders
# -------------------------------------------------------

def _build_title_slide(prs, course_name, module_title, lesson_title, module_number, lesson_number):
    """First slide: full navy background, course → module → lesson hierarchy."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

    # Full navy background
    _add_filled_rect(slide, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, BRAND_NAVY)

    # Accent stripe at the top
    _add_filled_rect(slide, 0, 0, SLIDE_WIDTH, Inches(0.08), BRAND_CYAN)

    # Logo
    _add_logo(slide)

    # Course name — small, upper area
    _add_text_box(
        slide,
        Inches(0.8), Inches(1.8), Inches(10), Inches(0.5),
        course_name.upper(), font_size=14, font_color=RGBColor(0xAA, 0xBB, 0xDD),
        bold=True, font_name="Calibri",
    )

    # Module label
    _add_text_box(
        slide,
        Inches(0.8), Inches(2.5), Inches(10), Inches(0.5),
        f"Module {module_number}  ·  {module_title}",
        font_size=18, font_color=BRAND_CYAN, bold=False,
    )

    # Lesson title — the hero text
    _add_text_box(
        slide,
        Inches(0.8), Inches(3.3), Inches(11), Inches(2.0),
        lesson_title,
        font_size=36, font_color=WHITE, bold=True,
    )

    # Bottom accent bar
    _add_filled_rect(
        slide,
        Inches(0.8), Inches(5.8), Inches(3), Inches(0.06), BRAND_ORANGE
    )

    # Lesson number badge
    _add_text_box(
        slide,
        Inches(0.8), Inches(6.2), Inches(4), Inches(0.4),
        f"Lesson {lesson_number}",
        font_size=14, font_color=BRAND_ORANGE, bold=True,
    )


def _build_content_slide(prs, slide_title, bullets, slide_num, total_slides, script=None, images=None):
    """Standard content slide: navy header bar + white body with bullets."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

    # Top accent line
    _add_filled_rect(slide, 0, 0, SLIDE_WIDTH, Inches(0.06), BRAND_CYAN)

    # Navy header bar
    _add_filled_rect(slide, 0, Inches(0.06), SLIDE_WIDTH, HEADER_HEIGHT, BRAND_NAVY)

    # Slide title in the header
    _add_text_box(
        slide,
        Inches(0.8), Inches(0.3), Inches(11), Inches(0.9),
        slide_title, font_size=28, font_color=WHITE, bold=True,
    )

    # Logo in header
    _add_logo(slide)

    # Light background for content area
    _add_filled_rect(
        slide,
        Inches(0.4), CONTENT_TOP,
        SLIDE_WIDTH - Inches(0.8), CONTENT_HEIGHT,
        LIGHT_GRAY,
    )

    # Bullets
    if bullets:
        # If slide has images, use a narrower width for the bullet text box (split layout)
        width = Inches(7.0) if images else CONTENT_WIDTH
        txBox = slide.shapes.add_textbox(
            CONTENT_LEFT, Inches(1.9), width, CONTENT_HEIGHT - Inches(0.3)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.auto_size = None

        for i, bullet_text in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            p.alignment = PP_ALIGN.LEFT
            p.space_before = Pt(10)
            p.space_after = Pt(6)
            p.level = 0

            # Bullet marker (blue dot via text)
            marker_run = p.add_run()
            marker_run.text = "●  "
            marker_run.font.size = Pt(12)
            marker_run.font.color.rgb = BRAND_NAVY
            marker_run.font.name = "Calibri"

            # Bullet text
            text_run = p.add_run()
            text_run.text = bullet_text
            text_run.font.size = Pt(18)
            text_run.font.color.rgb = BULLET_GRAY
            text_run.font.name = "Calibri"
            text_run.font.bold = False

    # Add image if present (bullets left, image right split layout)
    if images:
        img_data = images[0]
        img_path = img_data.get("file_path")
        abs_img_path = os.path.join(BASE_DIR, img_path)
        
        if os.path.exists(abs_img_path):
            try:
                # Add the picture initially at native size to get dimensions
                RIGHT_CONTENT_LEFT = CONTENT_LEFT + Inches(7.0) + Inches(0.4)
                pic = slide.shapes.add_picture(abs_img_path, RIGHT_CONTENT_LEFT, Inches(1.9))
                
                max_w = Inches(4.3)
                max_h = Inches(4.8)
                native_w = pic.width
                native_h = pic.height
                
                # Compute proportional scaling factors to fit within max bounds
                scale = min(max_w / native_w, max_h / native_h)
                pic.width = int(native_w * scale)
                pic.height = int(native_h * scale)
                
                # Center the scaled image vertically and horizontally in the right-hand area
                pic.top = Inches(1.9) + int((max_h - pic.height) / 2)
                pic.left = RIGHT_CONTENT_LEFT + int((max_w - pic.width) / 2)
            except Exception as e:
                print(f"  [SLIDES][WARNING] Failed to insert image {img_path}: {e}")

    # Slide number
    _add_slide_number(slide, slide_num, total_slides)

    # Bottom accent bar
    _add_filled_rect(slide, 0, SLIDE_HEIGHT - Inches(0.06), SLIDE_WIDTH, Inches(0.06), BRAND_NAVY)

    # Speaker notes script
    if script:
        try:
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = script
        except Exception as e:
            print(f"  [SLIDES][WARNING] Failed to write speaker notes: {e}")


def _build_end_slide(prs, lesson_title, module_title):
    """Final slide: navy background, end-of-lesson marker."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Full navy
    _add_filled_rect(slide, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, BRAND_NAVY)

    # Accent stripe
    _add_filled_rect(slide, 0, 0, SLIDE_WIDTH, Inches(0.08), BRAND_ORANGE)

    _add_logo(slide)

    # Completion text
    _add_text_box(
        slide,
        Inches(0.8), Inches(2.5), Inches(10), Inches(0.6),
        "LESSON COMPLETE", font_size=16, font_color=BRAND_ORANGE, bold=True,
    )

    # Lesson title recap
    _add_text_box(
        slide,
        Inches(0.8), Inches(3.3), Inches(10), Inches(1.5),
        lesson_title, font_size=32, font_color=WHITE, bold=True,
    )

    # Module reference
    _add_text_box(
        slide,
        Inches(0.8), Inches(5.2), Inches(10), Inches(0.5),
        module_title, font_size=16, font_color=BRAND_CYAN, bold=False,
    )

    # Bottom bar
    _add_filled_rect(
        slide, Inches(0.8), Inches(6.0), Inches(3), Inches(0.06), BRAND_CYAN
    )


# -------------------------------------------------------
# Main Generator
# -------------------------------------------------------

def generate_lesson_pptx(
    course: dict,
    module_index: int,
    lesson_index: int,
) -> str:
    """
    Generate a branded .pptx for one lesson.

    Returns the absolute path of the generated file.
    """
    course_id = course.get("id", "unknown")
    course_name = course.get("course_name", "Untitled Course")
    modules = course.get("modules", [])

    if module_index >= len(modules):
        raise ValueError(f"Module index {module_index} out of range (course has {len(modules)} modules).")

    module = modules[module_index]
    module_title = module.get("title", f"Module {module_index + 1}")
    module_number = module.get("module_number", module_index + 1)
    lessons = module.get("lessons", [])

    if lesson_index >= len(lessons):
        raise ValueError(f"Lesson index {lesson_index} out of range (module has {len(lessons)} lessons).")

    lesson = lessons[lesson_index]
    lesson_title = lesson.get("lesson_title", f"Lesson {lesson_index + 1}")
    lesson_number = lesson.get("lesson_number", lesson_index + 1)
    slides = lesson.get("slides", [])

    # ---- Build the presentation ----
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    total_content_slides = len(slides)
    total_slides_with_bookends = total_content_slides + 2  # title + end

    # Title slide
    _build_title_slide(prs, course_name, module_title, lesson_title, module_number, lesson_number)

    # Content slides
    for i, slide_data in enumerate(slides):
        s_title = slide_data.get("slide_title", f"Slide {i + 1}")
        bullet_texts = [b.get("text", "") for b in slide_data.get("bullets", []) if b.get("text", "").strip()]
        script = slide_data.get("script", "")
        _build_content_slide(prs, s_title, bullet_texts, i + 1, total_content_slides, script, slide_data.get("images", []))

    # End slide
    _build_end_slide(prs, lesson_title, module_title)

    # ---- Save ----
    output_dir = os.path.join(SLIDES_DIR, course_id)
    os.makedirs(output_dir, exist_ok=True)

    filename = f"module_{module_index + 1}_lesson_{lesson_index + 1}.pptx"
    output_path = os.path.join(output_dir, filename)
    prs.save(output_path)

    print(f"  [SLIDES] Generated {output_path} ({total_content_slides} content slides)")
    return output_path


def generate_all_slides_for_course(course_id: str) -> dict:
    """
    Generate .pptx for every lesson in every module of a course.
    Returns a manifest dict: {module_index: {lesson_index: filepath}}.
    """
    if not os.path.exists(COURSES_FILE):
        raise FileNotFoundError("Courses database not found.")

    with open(COURSES_FILE, "r", encoding="utf-8") as f:
        courses = json.load(f)

    course = next((c for c in courses if c.get("id") == course_id), None)
    if course is None:
        raise ValueError(f"Course '{course_id}' not found.")

    modules = course.get("modules", [])
    manifest = {}

    for mi, module in enumerate(modules):
        lessons = module.get("lessons", [])
        if not lessons:
            continue
        manifest[mi] = {}
        for li, lesson in enumerate(lessons):
            slides = lesson.get("slides", [])
            if not slides:
                continue
            try:
                path = generate_lesson_pptx(course, mi, li)
                manifest[mi][li] = path
            except Exception as e:
                print(f"  [SLIDES][ERROR] Module {mi+1} Lesson {li+1}: {e}")
                manifest[mi][li] = None

    total = sum(1 for m in manifest.values() for p in m.values() if p)
    print(f"  [SLIDES] Generated {total} PPTX files for course '{course.get('course_name')}'.")
    return manifest


def get_slide_path(course_id: str, module_index: int, lesson_index: int) -> str | None:
    """Return the expected path for a lesson PPTX, or None if not found."""
    filename = f"module_{module_index + 1}_lesson_{lesson_index + 1}.pptx"
    path = os.path.join(SLIDES_DIR, course_id, filename)
    return path if os.path.exists(path) else None


def list_available_slides(course_id: str) -> list[dict]:
    """Return a list of all generated slide files for a course."""
    course_dir = os.path.join(SLIDES_DIR, course_id)
    if not os.path.exists(course_dir):
        return []

    results = []
    for filename in sorted(os.listdir(course_dir)):
        if filename.endswith(".pptx"):
            # Parse module_X_lesson_Y.pptx
            parts = filename.replace(".pptx", "").split("_")
            try:
                mi = int(parts[1]) - 1  # 0-indexed
                li = int(parts[3]) - 1
                results.append({
                    "module_index": mi,
                    "lesson_index": li,
                    "filename": filename,
                    "path": os.path.join(course_dir, filename),
                })
            except (IndexError, ValueError):
                continue

    return results
