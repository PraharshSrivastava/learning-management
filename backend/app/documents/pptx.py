"""Native PPTX text extraction and LLM metadata inference."""

from __future__ import annotations

import re
from typing import List

from app.core.logging import generation_logger
from app.core.providers import get_llm_endpoint, safe_chat_completion
from app.generation.prompts import load_prompt
from app.generation.runtime import retry
from app.schemas.generation.blueprint import CourseMetadataSchema

logger = generation_logger(__name__)


def _clean_text(text: str) -> str:
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    text = re.sub(r"[ \t]+", " ", text)
    text = "\n".join(line.strip() for line in text.split("\n"))
    return re.sub(r"\n{3,}", "\n\n", text).strip()


def _sentence_lines(text: str) -> str:
    sentence_split = re.compile(r"(?<=[.!?])\s+")
    lines: list[str] = []
    for raw_line in text.split("\n"):
        stripped = raw_line.strip()
        if not stripped:
            continue
        lines.extend(part.strip() for part in sentence_split.split(stripped) if part.strip())
    return "\n".join(lines)


def extract_text_from_pptx(pptx_path: str) -> List[str]:
    """Extract and normalize visible text from PPTX text frames and tables."""
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError(
            "PPTX support requires python-pptx. Install backend requirements and restart."
        ) from exc

    logger.info("pptx_text_extraction_started path=%s", pptx_path)
    presentation = Presentation(pptx_path)
    body_lines: List[str] = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        slide_text: List[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                for paragraph in shape.text_frame.paragraphs:
                    paragraph_text = paragraph.text.strip()
                    if paragraph_text:
                        slide_text.append(paragraph_text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    cells = [
                        re.sub(r"\s+", " ", cell.text).strip()
                        for cell in row.cells
                        if cell.text and cell.text.strip()
                    ]
                    if cells:
                        slide_text.append(" | ".join(cells))

        cleaned_text = _clean_text("\n".join(slide_text))
        normalized_text = _sentence_lines(cleaned_text)
        slide_lines = [line.strip() for line in normalized_text.split("\n") if line.strip()]
        body_lines.extend(slide_lines)
        logger.info(
            "pptx_slide_text_extracted slide=%s line_count=%s",
            slide_number,
            len(slide_lines),
        )

    logger.info(
        "pptx_text_extraction_completed slide_count=%s line_count=%s char_count=%s",
        len(presentation.slides),
        len(body_lines),
        len("\n".join(body_lines)),
    )
    return body_lines


def extract_pptx_metadata_with_llm(body_lines: List[str], course_id: str = "blueprint") -> dict:
    numbered_text = "\n".join(
        f"[LINE {line_number}] {line}" for line_number, line in enumerate(body_lines, start=1)
    )
    json_schema = CourseMetadataSchema.model_json_schema()

    def generate_once():
        base_url, model_name = get_llm_endpoint(purpose="metadata")
        response = safe_chat_completion(
            base_url=base_url,
            model=model_name,
            messages=[
                {"role": "system", "content": load_prompt("pptx_metadata_extraction.md")},
                {
                    "role": "user",
                    "content": (
                        "Infer course metadata from the cleaned PowerPoint text below.\n\n"
                        f"{numbered_text}"
                    ),
                },
            ],
            response_format={
                "type": "json_schema",
                "json_schema": {
                    "name": "CourseMetadataSchema",
                    "schema": json_schema,
                },
            },
            temperature=0.1,
            default_max_tokens=1024,
            course_id=course_id,
            stage="pptx_metadata",
            attempts=1,
        )
        parsed = CourseMetadataSchema.model_validate_json(response.choices[0].message.content)
        return parsed.model_dump()

    return retry(generate_once, course_id=course_id, stage="pptx_metadata", attempts=3)
